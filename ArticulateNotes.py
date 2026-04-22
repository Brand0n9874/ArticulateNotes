import os
import shutil
import zipfile
import subprocess
import re
from pathlib import Path

import pandas as pd
import tkinter as tk
from tkinter import filedialog, messagebox

from pptx import Presentation
from openpyxl import load_workbook
from openpyxl.styles import Font, PatternFill
from openpyxl.worksheet.hyperlink import Hyperlink
from openpyxl.utils import get_column_letter


SCRIPT_VERSION = "6.0-ONE-SHEET-PPTX_NOTES-AUDIOLINKS-COL-D"
AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".aac", ".wma"}

# Keep extracted PPTA folders so hyperlinks stay valid on disk
DEFAULT_KEEP_PPTA_EXTRACTED = True

# Excel/OpenXML cannot contain certain control characters (often appear in notes)
_ILLEGAL_XML_RE = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F]")


# ----------------------------
# GUI
# ----------------------------
def launch_gui():
    selected_path = {"value": None}

    root = tk.Tk()
    root.title("Articulate Extractor (PPTX Notes + Audio Links)")
    root.geometry("760x380")
    root.resizable(False, False)
    root.configure(bg="#f0f4f8")

    tk.Label(
        root,
        text="Articulate Extractor",
        font=("Segoe UI", 18, "bold"),
        bg="#f0f4f8",
        fg="#1a1a2e",
    ).pack(pady=(22, 6))

    tk.Label(
        root,
        text=(
            "Output:\n"
            "• ONE Excel sheet: PPTX_Notes\n"
            "• Column D = Audio Link (hyperlink)\n"
            "• Notes column moved to the right\n\n"
            "Audio links are assigned in best-effort order (Slide 1 → Audio 1, etc.)."
        ),
        font=("Segoe UI", 10),
        bg="#f0f4f8",
        fg="#555555",
        justify="center",
    ).pack(pady=(0, 14))

    path_var = tk.StringVar(value="No file or folder selected")
    tk.Label(
        root,
        textvariable=path_var,
        font=("Segoe UI", 9),
        bg="#e8edf2",
        fg="#333333",
        relief="sunken",
        anchor="w",
        padx=8,
        width=90,
        height=2,
        wraplength=720,
        justify="left",
    ).pack(pady=(0, 12), padx=18)

    options_frame = tk.Frame(root, bg="#f0f4f8")
    options_frame.pack(pady=(0, 8))

    keep_var = tk.BooleanVar(value=DEFAULT_KEEP_PPTA_EXTRACTED)
    tk.Checkbutton(
        options_frame,
        text="Keep extracted PPTA folders (recommended so Excel audio links work)",
        variable=keep_var,
        bg="#f0f4f8",
        fg="#333333",
        font=("Segoe UI", 9),
        activebackground="#f0f4f8",
    ).pack()

    btn_frame = tk.Frame(root, bg="#f0f4f8")
    btn_frame.pack(pady=(8, 0))

    run_btn = tk.Button(
        root,
        text="Run Extraction",
        font=("Segoe UI", 11, "bold"),
        bg="#00b050",
        fg="white",
        activebackground="#007a35",
        activeforeground="white",
        width=22,
        height=2,
        bd=0,
        cursor="hand2",
        state="disabled",
    )

    def select_file():
        path = filedialog.askopenfilename(
            title="Select a ZIP, PPTA, or PPTX file",
            filetypes=[
                ("Supported files", "*.zip *.ppta *.pptx"),
                ("All files", "*.*"),
            ],
        )
        if path:
            selected_path["value"] = path
            path_var.set(f"FILE: {path}")
            run_btn.config(state="normal")

    def select_folder():
        path = filedialog.askdirectory(title="Select a folder")
        if path:
            selected_path["value"] = path
            path_var.set(f"FOLDER: {path}")
            run_btn.config(state="normal")

    btn_style = {
        "font": ("Segoe UI", 10, "bold"),
        "width": 18,
        "height": 2,
        "bd": 0,
        "cursor": "hand2",
        "fg": "white",
        "bg": "#0077b6",
        "activebackground": "#005f8e",
        "activeforeground": "white",
    }

    tk.Button(btn_frame, text="Select File", command=select_file, **btn_style).grid(row=0, column=0, padx=10)
    tk.Button(btn_frame, text="Select Folder", command=select_folder, **btn_style).grid(row=0, column=1, padx=10)

    def on_run():
        root.keep_extracted = keep_var.get()
        root.destroy()

    run_btn.config(command=on_run)
    run_btn.pack(pady=(18, 0))
    root.mainloop()

    return selected_path["value"], getattr(root, "keep_extracted", DEFAULT_KEEP_PPTA_EXTRACTED)


# ----------------------------
# Utilities
# ----------------------------
def unzip_if_needed(inpath: Path, work_dir: Path) -> Path:
    """Folder -> folder; Zip -> extract then folder; File -> parent folder."""
    if inpath.is_dir():
        return inpath

    if zipfile.is_zipfile(str(inpath)):
        target = work_dir / "unzipped"
        if target.exists():
            shutil.rmtree(target, ignore_errors=True)
        target.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(str(inpath), "r") as z:
            z.extractall(str(target))
        return target

    return inpath.parent


def find_7zip() -> str:
    candidates = [
        r"C:\Program Files\7-Zip\7z.exe",
        r"C:\Program Files (x86)\7-Zip\7z.exe",
    ]
    for c in candidates:
        if os.path.isfile(c):
            return c
    found = shutil.which("7z") or shutil.which("7za")
    if found:
        return found
    raise FileNotFoundError(
        "7-Zip not found. Install 7-Zip or ensure '7z' is on your PATH.\n"
        "Download: https://www.7-zip.org/"
    )


def _sanitize_excel_text(value):
    """Remove illegal XML control chars that can corrupt XLSX and trigger Excel repair."""
    if value is None:
        return None
    if isinstance(value, str):
        return _ILLEGAL_XML_RE.sub("", value)
    return value


def _sanitize_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    """Sanitize every cell for Excel safety."""
    if hasattr(df, "map"):  # pandas >= 2.x
        return df.map(_sanitize_excel_text)
    return df.applymap(_sanitize_excel_text)


def _path_to_file_uri(path_str: str) -> str:
    """Correct Windows-safe file URI (prevents Excel 'repair' prompts)."""
    return Path(path_str).resolve().as_uri()


def autosize_columns(ws, max_width=88):
    for col in range(1, ws.max_column + 1):
        letter = get_column_letter(col)
        max_len = 10
        for row in range(1, ws.max_row + 1):
            v = ws.cell(row=row, column=col).value
            if v is None:
                continue
            max_len = max(max_len, len(str(v)))
        ws.column_dimensions[letter].width = min(max_len + 3, max_width)


# ----------------------------
# PPTX: Notes extraction
# ----------------------------
def extract_pptx_notes(pptx_path: Path) -> list[dict]:
    rows = []
    try:
        prs = Presentation(str(pptx_path))
        for idx, slide in enumerate(prs.slides, start=1):
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()

            rows.append(
                {
                    "PPTX File": pptx_path.name,
                    "Slide Number": idx,
                    "Notes Found": bool(notes_text),
                    # Audio Link is filled later via openpyxl (must be column D)
                    "Audio Link": "",
                    "Notes": _sanitize_excel_text(notes_text) if notes_text else "NO NOTES",
                }
            )
    except Exception as e:
        rows.append(
            {
                "PPTX File": pptx_path.name,
                "Slide Number": "ERROR",
                "Notes Found": False,
                "Audio Link": "",
                "Notes": _sanitize_excel_text(f"ERROR reading PPTX: {e}"),
            }
        )
    return rows


def get_pptx_slide_count(pptx_path: Path) -> int:
    try:
        prs = Presentation(str(pptx_path))
        return len(prs.slides)
    except Exception:
        return 0


# ----------------------------
# PPTA: Extract audio files (no output sheet)
# ----------------------------
def extract_ppta_to_folder(ppta_path: Path, extract_root: Path) -> Path:
    """Extract PPTA using 7-Zip into extract_root/<ppta_stem>/ and return extraction folder."""
    seven_zip = find_7zip()
    target = extract_root / ppta_path.stem
    target.mkdir(parents=True, exist_ok=True)

    result = subprocess.run(
        [seven_zip, "x", str(ppta_path), f"-o{str(target)}", "-y"],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        encoding="utf-8",
        errors="replace",
    )

    if result.returncode != 0:
        raise RuntimeError(
            f"7-Zip failed extracting: {ppta_path}\n"
            f"Return code: {result.returncode}\n"
            f"STDOUT:\n{result.stdout}\n\nSTDERR:\n{result.stderr}"
        )

    return target


def collect_audio_paths(extracted_folder: Path) -> list[str]:
    """
    Return a STABLE ordered list of full audio paths from the extracted PPTA.
    Sorting by relative path + filename gives consistent ordering run-to-run.
    """
    items = []
    for f in extracted_folder.rglob("*"):
        if f.is_file() and f.suffix.lower() in AUDIO_EXTS:
            rel = str(f.relative_to(extracted_folder)).replace("\\", "/").lower()
            items.append((rel, f.name.lower(), str(f.resolve())))
    items.sort()
    return [p for _, __, p in items]


def choose_partner_ppta_for_pptx(pptx_path: Path, ppta_files: list[Path]) -> Path | None:
    """
    Best-effort pairing:
      - if only one PPTA exists, use it
      - else match by stem (Project.pptx -> Project.ppta)
    """
    if not ppta_files:
        return None
    if len(ppta_files) == 1:
        return ppta_files[0]
    match = next((p for p in ppta_files if p.stem.lower() == pptx_path.stem.lower()), None)
    return match


def build_assignments(pptx_files: list[Path], ppta_files: list[Path], audio_lists_by_ppta: dict[str, list[str]]) -> dict:
    """
    Return mapping:
      (pptx_filename, slide_number) -> audio_full_path

    Assignment strategy (for now):
      Slide 1 -> audio[0], Slide 2 -> audio[1], etc. (per partner PPTA).
    """
    assignments = {}
    for pptx in pptx_files:
        partner = choose_partner_ppta_for_pptx(pptx, ppta_files)
        audio_list = audio_lists_by_ppta.get(partner.name, []) if partner else []
        slide_count = get_pptx_slide_count(pptx)

        for s in range(1, slide_count + 1):
            audio_path = audio_list[s - 1] if (s - 1) < len(audio_list) else ""
            assignments[(pptx.name, s)] = audio_path
    return assignments


# ----------------------------
# Excel: Write ONE sheet + add hyperlinks in col D
# ----------------------------
def write_one_sheet_excel(out_path: Path, pptx_rows: list[dict], assignments: dict):
    """
    Writes ONE sheet only: PPTX_Notes
    Columns:
      A PPTX File
      B Slide Number
      C Notes Found
      D Audio Link  (hyperlink text)
      E Notes
    """
    df = pd.DataFrame(pptx_rows) if pptx_rows else pd.DataFrame(
        columns=["PPTX File", "Slide Number", "Notes Found", "Audio Link", "Notes"]
    )

    # Force exact order so column D is Audio Link
    for col in ["PPTX File", "Slide Number", "Notes Found", "Audio Link", "Notes"]:
        if col not in df.columns:
            df[col] = ""
    df = df[["PPTX File", "Slide Number", "Notes Found", "Audio Link", "Notes"]]
    df = _sanitize_dataframe(df)

    with pd.ExcelWriter(str(out_path), engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="PPTX_Notes")

    # Post-process: add hyperlinks + formatting
    wb = load_workbook(str(out_path))
    ws = wb["PPTX_Notes"]

    hyperlink_font = Font(color="0000FF", underline="single")
    missing_font = Font(color="FF0000")
    green_fill = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
    red_fill = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")

    # Identify columns (should be fixed, but stay robust)
    headers = [c.value for c in ws[1]]
    col_pptx = headers.index("PPTX File") + 1
    col_slide = headers.index("Slide Number") + 1
    col_nf = headers.index("Notes Found") + 1
    col_audio = headers.index("Audio Link") + 1  # should be 4
    # Notes column exists but we don't need its index here

    for r in range(2, ws.max_row + 1):
        pptx_name = ws.cell(row=r, column=col_pptx).value
        slide_num = ws.cell(row=r, column=col_slide).value

        # Color Notes Found
        nf_cell = ws.cell(row=r, column=col_nf)
        if nf_cell.value is True:
            nf_cell.fill = green_fill
        elif nf_cell.value is False:
            nf_cell.fill = red_fill

        # Add hyperlink in column D
        link_cell = ws.cell(row=r, column=col_audio)
        link_cell.value = ""
        link_cell.hyperlink = None

        if pptx_name and isinstance(slide_num, int):
            audio_path = assignments.get((str(pptx_name), slide_num), "")
            if audio_path and os.path.isfile(str(audio_path)):
                uri = _path_to_file_uri(str(audio_path))
                link_cell.value = "Open Audio"
                link_cell.hyperlink = Hyperlink(ref=link_cell.coordinate, target=uri)
                link_cell.font = hyperlink_font
            elif audio_path:
                link_cell.value = "Path Not Found"
                link_cell.font = missing_font

    ws.auto_filter.ref = ws.dimensions
    autosize_columns(ws)
    wb.save(str(out_path))


# ----------------------------
# Main
# ----------------------------
def main():
    inpath_str, keep_extracted = launch_gui()
    if not inpath_str:
        messagebox.showinfo("Cancelled", "No input selected.")
        return

    inpath = Path(inpath_str)
    base_dir = inpath if inpath.is_dir() else inpath.parent

    work_dir = base_dir / "_articulate_work"
    work_dir.mkdir(parents=True, exist_ok=True)

    extract_root = base_dir / "_ppta_extracted"
    extract_root.mkdir(parents=True, exist_ok=True)

    try:
        scan_root = unzip_if_needed(inpath, work_dir)

        pptx_files = sorted(scan_root.rglob("*.pptx"))
        ppta_files = sorted(scan_root.rglob("*.ppta"))

        # If user selected a single file, narrow scope
        if inpath.is_file():
            if inpath.suffix.lower() == ".pptx":
                pptx_files = [inpath]
                ppta_files = []
            elif inpath.suffix.lower() == ".ppta":
                ppta_files = [inpath]
                pptx_files = []

        # 1) Collect PPTX notes rows
        pptx_rows = []
        for pptx in pptx_files:
            pptx_rows.extend(extract_pptx_notes(pptx))

        # 2) Extract PPTA(s) and build audio lists (not written as sheets)
        audio_lists_by_ppta = {}
        for ppta in ppta_files:
            extracted = extract_ppta_to_folder(ppta, extract_root)
            audio_lists_by_ppta[ppta.name] = collect_audio_paths(extracted)

        if not pptx_rows:
            messagebox.showinfo("No Data", "No PPTX files found in the selected input.")
            return

        # 3) Build best-effort slide->audio assignment
        assignments = build_assignments(pptx_files, ppta_files, audio_lists_by_ppta)

        # 4) Write ONE-sheet Excel with audio hyperlinks in column D
        out_name = inpath.stem if inpath.is_file() else inpath.name
        out_path = base_dir / f"{out_name}_PPTX_notes.xlsx"
        write_one_sheet_excel(out_path, pptx_rows, assignments)

        # Cleanup: temp unzip folder; extracted audio depends on checkbox
        if work_dir.exists():
            shutil.rmtree(work_dir, ignore_errors=True)
        if not keep_extracted and extract_root.exists():
            shutil.rmtree(extract_root, ignore_errors=True)

        messagebox.showinfo(
            "Success",
            "Extraction complete!\n\n"
            f"Excel saved to:\n{out_path}\n\n"
            "Output workbook contains ONE sheet only (PPTX_Notes).\n"
            "Column D contains the audio hyperlinks; Notes column is shifted right.",
        )

    except Exception as e:
        # Clean temp unzip folder on error; keep extracted_root for debugging unless user unchecked
        if work_dir.exists():
            shutil.rmtree(work_dir, ignore_errors=True)
        messagebox.showerror("Error", str(e))
        raise


if __name__ == "__main__":
    print(f"Running version: {SCRIPT_VERSION}")
    main()